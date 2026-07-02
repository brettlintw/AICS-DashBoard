import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import os
import plotly.io as pio
from contextlib import contextmanager

# 頁面設定
st.set_page_config(layout="wide", page_title="AICS 北美決策中心")
st.title("🌐 AICS 北美部署決策中心 (V1.0 版)")

# 分享連結用的共用資料存放位置：管理者上傳的檔案會存在這裡，訪客（viewer）從這裡讀取，
# 不會出現在 git 版本紀錄裡（見 .gitignore），Streamlit Cloud 重新部署時會被清空，須重新上傳。
SHARED_DATA_PATH = "data/shared_dataset.xlsx"


def check_access():
    if st.session_state.get("role"):
        return st.session_state["role"]

    try:
        admin_pw = st.secrets.get("admin_password")
        viewer_pw = st.secrets.get("viewer_password")
    except Exception:
        # 尚未設定 .streamlit/secrets.toml（本機開發或還沒設定 Streamlit Cloud Secrets）時會拋例外。
        admin_pw = viewer_pw = None
    if not admin_pw and not viewer_pw:
        # 尚未在 Streamlit Cloud 的 Secrets 設定密碼時，維持原本不需要密碼即可使用的行為。
        st.session_state["role"] = "admin"
        return "admin"

    st.subheader("🔒 請輸入存取密碼")
    pwd = st.text_input("密碼", type="password")
    if not pwd:
        st.stop()
    if admin_pw and pwd == admin_pw:
        st.session_state["role"] = "admin"
    elif viewer_pw and pwd == viewer_pw:
        st.session_state["role"] = "viewer"
    else:
        st.error("密碼錯誤")
        st.stop()
    st.rerun()


role = check_access()

# 北美州代碼與座標映射表
US_STATES_COORDS = {'AL': [32.8, -86.7], 'AK': [61.3, -152.4], 'AZ': [33.7, -111.4], 'AR': [34.9, -92.3], 'CA': [36.1, -119.6], 'CO': [39.0, -105.3], 'CT': [41.5, -72.7], 'DE': [39.3, -75.5], 'FL': [27.7, -81.6], 'GA': [33.0, -83.6], 'HI': [21.0, -157.4], 'ID': [44.2, -114.4], 'IL': [40.3, -88.9], 'IN': [39.8, -86.2], 'IA': [42.0, -93.2], 'KS': [38.5, -96.7], 'KY': [37.6, -84.6], 'LA': [31.1, -91.8], 'ME': [44.6, -69.3], 'MD': [39.0, -76.8], 'MA': [42.2, -71.5], 'MI': [43.3, -84.5], 'MN': [45.6, -93.9], 'MS': [32.7, -89.6], 'MO': [38.4, -92.2], 'MT': [46.9, -110.4], 'NE': [41.1, -98.2], 'NV': [38.3, -117.0], 'NH': [43.4, -71.5], 'NJ': [40.2, -74.5], 'NM': [34.8, -106.2], 'NY': [42.1, -74.9], 'NC': [35.6, -79.8], 'ND': [47.5, -99.7], 'OH': [40.3, -82.7], 'OK': [35.5, -96.9], 'OR': [44.5, -122.0], 'PA': [40.5, -77.2], 'RI': [41.6, -71.5], 'SC': [33.8, -80.9], 'SD': [44.2, -99.4], 'TN': [35.7, -86.6], 'TX': [31.0, -97.5], 'UT': [40.1, -111.8], 'VT': [44.0, -72.7], 'VA': [37.7, -78.1], 'WA': [47.4, -120.4], 'WV': [38.4, -80.9], 'WI': [44.2, -89.6], 'WY': [42.7, -107.3]}

# 高對比氣泡配色（避免預設色系太相近分不清楚）
BUBBLE_COLORS = ['#e6194B', '#3cb44b', '#4363d8', '#f58231', '#911eb4', '#42d4f4', '#f032e6', '#469990', '#9A6324', '#000075']


def get_machine_type_colors(f_df):
    # 跟地圖氣泡用同一套配色規則，讓排名長條圖跟地圖上的顏色能對得起來。
    types = f_df['Machine Type'].unique()
    return {m: BUBBLE_COLORS[i % len(BUBBLE_COLORS)] for i, m in enumerate(types)}


@contextmanager
def export_color_theme():
    # Streamlit 會把圖表顏色暫時換成佔位色（如 #000001），交給瀏覽器顯示時才轉成真正顏色；
    # PPT 匯出繞過瀏覽器直接轉圖片，所以要暫時切回標準色版，轉完再切回來，避免畫面上的圖表跟著變色。
    original = pio.templates.default
    pio.templates.default = "plotly"
    try:
        yield
    finally:
        pio.templates.default = original


def build_map_fig(f_df):
    # 氣泡大小改用「目前篩選範圍內的最小~最大值」正規化到固定的可視像素區間，
    # 這樣不管篩選後數據是大是小（例如只選一天、或只剩個位數出貨量），氣泡都不會小到看不見。
    MIN_BUBBLE, MAX_BUBBLE = 12, 45
    state_totals = f_df.groupby(['Machine Type', 'State Code'])['Outbound Qty (Item)'].sum()
    qty_min = state_totals.min() if len(state_totals) else 0
    qty_max = state_totals.max() if len(state_totals) else 0

    def scaled_size(qty):
        if qty_max == qty_min:
            return MAX_BUBBLE
        return MIN_BUBBLE + (qty - qty_min) / (qty_max - qty_min) * (MAX_BUBBLE - MIN_BUBBLE)

    # 國旗效果直接畫在各州的實際圖形上（Choropleth），而不是疊一張背景圖片，
    # 這樣顏色保證只會出現在真正畫出州界的地方，不會跑到地圖以外的空白區域。
    CANTON_STATES = {'WA', 'OR', 'ID', 'MT', 'WY', 'ND', 'SD', 'MN'}
    STRIPE_RED = 'rgba(178,34,52,0.16)'
    CANTON_BLUE = 'rgba(60,59,110,0.18)'
    TRANSPARENT = 'rgba(255,255,255,0)'
    N_BANDS = 9
    lats = [v[0] for v in US_STATES_COORDS.values()]
    lat_min, lat_max = min(lats), max(lats)

    def flag_zone(state):
        if state in CANTON_STATES:
            return 2
        band = int((lat_max - US_STATES_COORDS[state][0]) / (lat_max - lat_min + 1e-9) * N_BANDS)
        return band % 2

    all_states = list(US_STATES_COORDS.keys())
    type_colors = get_machine_type_colors(f_df)
    fig_map = go.Figure()
    fig_map.add_trace(go.Choropleth(
        locations=all_states, locationmode="USA-states", z=[flag_zone(s) for s in all_states],
        zmin=0, zmax=2, showscale=False, hoverinfo='skip', marker_line_color='rgba(0,0,0,0)',
        colorscale=[[0, TRANSPARENT], [0.33, TRANSPARENT], [0.34, STRIPE_RED], [0.66, STRIPE_RED], [0.67, CANTON_BLUE], [1, CANTON_BLUE]],
    ))
    fig_map.add_trace(go.Scattergeo(lon=[US_STATES_COORDS[s][1] for s in US_STATES_COORDS], lat=[US_STATES_COORDS[s][0] for s in US_STATES_COORDS], text=list(US_STATES_COORDS.keys()), mode='text', textfont=dict(size=14, color='blue'), showlegend=False))
    for m in f_df['Machine Type'].unique():
        m_df = f_df[f_df['Machine Type'] == m].groupby('State Code')['Outbound Qty (Item)'].sum().reset_index()
        fig_map.add_trace(go.Scattergeo(locations=m_df['State Code'], locationmode="USA-states",
                                         marker=dict(size=m_df['Outbound Qty (Item)'].apply(scaled_size),
                                                      color=type_colors[m],
                                                      line=dict(width=1, color='white')),
                                         text=m_df['Outbound Qty (Item)'], hovertemplate="%{location}: %{text}<extra></extra>",
                                         name=m))
    fig_map.update_layout(
        geo=dict(scope='usa', projection=dict(type='albers usa'), landcolor='#fbfbfd', lakecolor='#fbfbfd', bgcolor='white', subunitcolor='#9aa5b1'),
        paper_bgcolor='white',
        height=600, margin={"l": 0, "r": 0, "t": 0, "b": 60}, legend=dict(orientation='h', x=0.5, xanchor='center', y=-0.1))
    return fig_map


def build_type_rank_fig(f_df):
    type_summary = f_df.groupby('Machine Type')['Outbound Qty (Item)'].sum()
    type_colors = get_machine_type_colors(f_df)
    fig = go.Figure(go.Bar(
        x=type_summary.values, y=type_summary.index, orientation='h',
        marker_color=[type_colors[m] for m in type_summary.index],
        text=type_summary.values, textposition='outside'))
    fig.update_layout(
        yaxis=dict(categoryorder='total ascending', title=None), xaxis=dict(title='出貨量'),
        height=280, margin={"l": 10, "r": 30, "t": 10, "b": 30}, showlegend=False)
    return fig


def build_dim_fig(df_g, dim, mode, chart_type):
    if chart_type == "柱狀圖":
        fig = px.bar(df_g, x=dim if mode == "全時間段" else "Month-Date", y='Outbound Qty (Item)', color=dim, text='Outbound Qty (Item)')
    elif chart_type == "推移圖":
        fig = px.line(df_g, x="Month-Date", y='Outbound Qty (Item)', color=dim, markers=True, text='Outbound Qty (Item)')
        fig.update_traces(textposition="top center")
        fig.update_xaxes(dtick="M1", tickformat="%Y-%m", tickangle=45)
    else:
        fig = px.pie(df_g, values='Outbound Qty (Item)', names=dim)
    return fig


def build_abc_fig(abc_df):
    fig_abc = go.Figure()
    fig_abc.add_trace(go.Bar(x=abc_df['Company'], y=abc_df['Outbound Qty (Item)'], name='出貨量',
                              marker_color=abc_df['分級'].map({'A': '#d62728', 'B': '#ff7f0e', 'C': '#1f77b4'})))
    fig_abc.add_trace(go.Scatter(x=abc_df['Company'], y=abc_df['累計佔比(%)'], name='累計佔比(%)',
                                  yaxis='y2', mode='lines+markers'))
    fig_abc.update_layout(
        yaxis=dict(title='出貨量'),
        yaxis2=dict(title='累計佔比(%)', overlaying='y', side='right', range=[0, 100]),
        height=500, xaxis_tickangle=45, legend=dict(x=1.05, y=1))
    return fig_abc


def build_churn_fig(churn_list, churn_threshold):
    fig_churn = px.bar(churn_list, x='Company', y='距今月數', color='距今月數',
                        color_continuous_scale='Reds', title=f"超過 {churn_threshold} 個月無出貨的客戶")
    fig_churn.update_xaxes(tickangle=45)
    return fig_churn


def find_layout_by_name(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    return None


def remove_all_slides(prs):
    # 只把 <p:sldId> 從清單拿掉，範本裡原本那張投影片對應的關聯(relationship)還留著，
    # 沒被任何東西參照，PowerPoint 開啟時會判定簡報內容有問題、要求修復。
    # 兩個都要清乾淨才是正確的移除方式。
    xml_slides = prs.slides._sldIdLst
    for slide_id in list(xml_slides):
        prs.part.drop_rel(slide_id.rId)
        xml_slides.remove(slide_id)


def make_report_presentation(bg_template_bytes):
    # 有上傳背景範本時，沿用範本的簡報尺寸與版面設計（母片/裝飾圖形），而不是重新做一個空白簡報，
    # 這樣背景風格才會自動套用；沒有範本時維持原本的大版面，讓表格有更多空間。
    if bg_template_bytes:
        prs = Presentation(io.BytesIO(bg_template_bytes))
        remove_all_slides(prs)
        layout = find_layout_by_name(prs, "只有標題") or prs.slide_masters[0].slide_layouts[0]

        title_ph = layout.placeholders[0] if len(layout.placeholders) else None
        if title_ph is not None and title_ph.top is not None:
            content_top = title_ph.top + (title_ph.height or 0) + Inches(0.15)
        else:
            content_top = Inches(0.9)

        cfg = dict(
            layout=layout, left=Inches(0.4), width=Inches(9.2), content_top=content_top,
            safe_bottom=int(prs.slide_height * 0.87), row_height=Inches(0.2), max_rows_cap=8,
            header_font=Pt(8), cell_font=Pt(7), note_font=Pt(6), chart_px=(1600, 320), map_px=(1600, 480),
        )
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(10)
        layout = prs.slide_layouts[6]
        cfg = dict(
            layout=layout, left=Inches(0.5), width=Inches(12), content_top=Inches(0.9),
            safe_bottom=int(prs.slide_height * 0.95), row_height=Inches(0.32), max_rows_cap=10,
            header_font=Pt(11), cell_font=Pt(10), note_font=Pt(8), chart_px=(1600, 450), map_px=(1600, 700),
        )
    return prs, cfg


def _add_table_to_slide(slide, cfg, table_df, anchor_top=None, anchor_bottom=None):
    """畫表格。anchor_top：表格緊跟在圖片下方，剩餘空間能放幾筆就放幾筆。
    anchor_bottom：表格固定貼齊這個高度（版面最下方），列數由資料量決定，回傳表格頂端位置給呼叫端用來計算圖片還能放多大。"""
    if table_df is None or len(table_df) == 0:
        return anchor_bottom

    if anchor_top is not None:
        available = cfg['safe_bottom'] - anchor_top - Inches(0.3)
        max_rows = max(1, min(cfg['max_rows_cap'], int(available / cfg['row_height']) - 1))
        display_df = table_df.head(max_rows)
        rows, cols = display_df.shape
        table_height = cfg['row_height'] * (rows + 1)
        table_top = anchor_top
        show_note = len(table_df) > max_rows
    else:
        max_rows = min(len(table_df), cfg['max_rows_cap'])
        display_df = table_df.head(max_rows)
        rows, cols = display_df.shape
        table_height = cfg['row_height'] * (rows + 1)
        table_top = anchor_bottom - table_height
        show_note = False  # 貼底模式是給列數固定很少的資料用，版面沒有預留註腳空間

    table = slide.shapes.add_table(rows + 1, cols, cfg['left'], table_top, cfg['width'], table_height).table
    for r_idx in range(rows + 1):
        table.rows[r_idx].height = cfg['row_height']
    for i, col in enumerate(display_df.columns):
        cell = table.cell(0, i)
        cell.text = str(col)
        cell.text_frame.paragraphs[0].font.size = cfg['header_font']
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r + 1, c)
            cell.text = str(display_df.iloc[r, c])
            cell.text_frame.paragraphs[0].font.size = cfg['cell_font']
    if show_note:
        # 固定貼在版面左下角（而非緊跟在表格後面），字體縮小成註腳大小，避免搶版面。
        note_top = cfg['safe_bottom'] - Inches(0.22)
        note = slide.shapes.add_textbox(Inches(0.25), note_top, cfg['width'], Inches(0.22))
        note.text_frame.text = f"僅顯示前 {max_rows} 筆，完整資料請見網頁畫面"
        note.text_frame.paragraphs[0].font.size = cfg['note_font']
    return table_top


PX_PER_INCH = 160  # 匯出圖片解析度，太低字會模糊、太高檔案會變大


def _emu_to_px(emu):
    return max(1, int(emu / 914400 * PX_PER_INCH))


def add_chart_slide(prs, cfg, title, fig, table_df, is_map=False):
    slide = prs.slides.add_slide(cfg['layout'])
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    else:
        slide.shapes.add_textbox(cfg['left'], Inches(0.2), cfg['width'], Inches(0.5)).text_frame.text = title

    content_bottom = cfg['content_top']
    try:
        px_w, px_h = cfg['map_px'] if is_map else cfg['chart_px']
        img_bytes = pio.to_image(fig, format="png", width=px_w, height=px_h)
        pic = slide.shapes.add_picture(io.BytesIO(img_bytes), cfg['left'], cfg['content_top'], width=cfg['width'])
        content_bottom = cfg['content_top'] + pic.height
    except Exception:
        pass

    _add_table_to_slide(slide, cfg, table_df, anchor_top=content_bottom + Inches(0.2))
    return slide


def add_dual_chart_slide(prs, cfg, title, left_fig, right_fig, table_df):
    slide = prs.slides.add_slide(cfg['layout'])
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    else:
        slide.shapes.add_textbox(cfg['left'], Inches(0.2), cfg['width'], Inches(0.5)).text_frame.text = title

    gap = Inches(0.2)
    left_width = int(cfg['width'] * 0.6)
    right_width = cfg['width'] - left_width - gap
    right_left = cfg['left'] + left_width + gap

    # 表格先貼齊版面底端算出佔用高度，剩下的垂直空間直接拿去請 Plotly 用同樣比例畫圖，
    # 保證圖片剛好填滿可用空間、不變形也不會溢出，比事後縮放更準確。
    table_top = _add_table_to_slide(slide, cfg, table_df, anchor_bottom=cfg['safe_bottom'])
    image_area_bottom = (table_top - Inches(0.2)) if table_top is not None else cfg['safe_bottom']
    max_img_height = image_area_bottom - cfg['content_top']

    for fig, box_left, box_width in (
        (left_fig, cfg['left'], left_width),
        (right_fig, right_left, right_width),
    ):
        try:
            img_bytes = pio.to_image(fig, format="png", width=_emu_to_px(box_width), height=_emu_to_px(max_img_height))
            slide.shapes.add_picture(io.BytesIO(img_bytes), box_left, cfg['content_top'], width=box_width, height=max_img_height)
        except Exception:
            pass

    return slide


if role == "admin":
    uploaded_file = st.sidebar.file_uploader("上傳 Excel（分享連結的訪客會看到這份資料）", type=["xlsx"])
    if uploaded_file:
        os.makedirs(os.path.dirname(SHARED_DATA_PATH), exist_ok=True)
        with open(SHARED_DATA_PATH, "wb") as f:
            f.write(uploaded_file.getvalue())
        st.sidebar.success("已更新分享連結訪客看到的資料")
    data_source = uploaded_file
else:
    st.sidebar.info("檢視模式：資料由管理者提供，無法上傳或下載原始檔案")
    data_source = SHARED_DATA_PATH if os.path.exists(SHARED_DATA_PATH) else None

if not data_source:
    if role == "viewer":
        st.info("尚未有可瀏覽的資料，請聯繫管理者上傳。")
else:
    df = pd.read_excel(data_source, sheet_name='Data Base')
    df.columns = df.columns.str.strip()
    df['Date(出庫)'] = pd.to_datetime(df['Date(出庫)'])

    st.sidebar.markdown("### 全域篩選")
    date_range = st.sidebar.date_input("日期篩選", [df['Date(出庫)'].min().date(), df['Date(出庫)'].max().date()])

    filters = {}
    dims = [('Machine Type', '設備維度'), ('Field', '場域維度'), ('Area', '區域維度'), ('Company', '客戶維度'), ('Device/Platform', '平台維度')]
    for col, label in dims:
        selected = st.sidebar.multiselect(f"{label}篩選 (留空即全選)", df[col].unique().tolist())
        filters[col] = selected if selected else df[col].unique().tolist()

    f_df = df[(df['Date(出庫)'].dt.date >= date_range[0]) & (df['Date(出庫)'].dt.date <= date_range[1])].copy()
    for col, selected in filters.items():
        f_df = f_df[f_df[col].isin(selected)]
    f_df['Month-Date'] = f_df['Date(出庫)'].dt.strftime('%Y-%m')

    # 地圖渲染與設備總數統計
    st.subheader("🗺️ 北美設備戰術分佈")
    type_summary = f_df.groupby('Machine Type')['Outbound Qty (Item)'].sum()
    cols = st.columns(len(type_summary) if len(type_summary) > 0 else 1)
    for i, (m_type, total) in enumerate(type_summary.items()):
        cols[i].metric(label=f"總 {m_type}", value=int(total))

    fig_map = build_map_fig(f_df)
    st.plotly_chart(fig_map, use_container_width=True)

    st.markdown("##### 📊 設備類型出貨量排名")
    fig_type_rank = build_type_rank_fig(f_df)
    st.plotly_chart(fig_type_rank, use_container_width=True)

    # 五大分析維度
    dim_exports = []
    for dim, name in dims:
        st.markdown("---")
        st.subheader(f"📈 {name}")
        c1, c2 = st.columns([1, 4])
        with c1:
            mode = st.radio("模式", ["月份推移", "全時間段"], key=f"m_{dim}")
            chart_type = st.radio("圖表", ["柱狀圖", "推移圖", "餅圖"], key=f"c_{dim}")
            sort = st.selectbox("排序", ["預設", "由大至小", "由小至大"], key=f"s_{dim}") if mode == "全時間段" else "預設"
            show_table = st.checkbox("顯示表格", value=True, key=f"t_{dim}")
        with c2:
            df_g = f_df.groupby(['Month-Date', dim] if mode=="月份推移" else dim)[['Outbound Qty (Item)']].sum().reset_index()
            if mode == "全時間段" and sort != "預設": df_g = df_g.sort_values(by='Outbound Qty (Item)', ascending=(sort=="由小至大"))

            fig = build_dim_fig(df_g, dim, mode, chart_type)

            st.plotly_chart(fig, use_container_width=True)
            if show_table: st.dataframe(df_g, use_container_width=True)
        dim_exports.append((name, dim, mode, chart_type, df_g))

    # 客戶經營洞察：ABC 分析 + 流失預警
    st.markdown("---")
    st.subheader("🎯 客戶經營洞察")

    dim_filtered_df = df.copy()
    for col, selected in filters.items():
        dim_filtered_df = dim_filtered_df[dim_filtered_df[col].isin(selected)]

    tab_abc, tab_churn = st.tabs(["📊 客戶 ABC 分析", "⚠️ 流失預警"])

    with tab_abc:
        abc_df = f_df.groupby('Company')['Outbound Qty (Item)'].sum().reset_index()
        abc_df = abc_df.sort_values('Outbound Qty (Item)', ascending=False).reset_index(drop=True)
        total_qty = abc_df['Outbound Qty (Item)'].sum()
        abc_df['累計佔比(%)'] = (abc_df['Outbound Qty (Item)'].cumsum() / total_qty * 100).round(1) if total_qty else 0
        abc_df['分級'] = abc_df['累計佔比(%)'].apply(lambda p: 'A' if p <= 80 else ('B' if p <= 95 else 'C'))

        grade_count = abc_df['分級'].value_counts().reindex(['A', 'B', 'C']).fillna(0).astype(int)
        m1, m2, m3 = st.columns(3)
        for col_m, grade in zip([m1, m2, m3], ['A', 'B', 'C']):
            cnt = grade_count.get(grade, 0)
            pct = cnt / len(abc_df) * 100 if len(abc_df) else 0
            col_m.metric(label=f"{grade} 級客戶", value=f"{cnt} 家", delta=f"佔客戶數 {pct:.1f}%")
        st.caption("A級：累計貢獻前80%出貨量的核心客戶；B級：80~95%；C級：長尾客戶")

        fig_abc = build_abc_fig(abc_df)
        st.plotly_chart(fig_abc, use_container_width=True)
        show_table_abc = st.checkbox("顯示表格", value=True, key="t_abc")
        if show_table_abc: st.dataframe(abc_df, use_container_width=True)

    with tab_churn:
        churn_threshold = st.slider("流失警戒門檻（月）", min_value=1, max_value=12, value=3)
        ref_date = df['Date(出庫)'].max()
        last_ship = dim_filtered_df.groupby('Company')['Date(出庫)'].max().reset_index()
        last_ship.columns = ['Company', '最後出貨日']
        last_ship['距今月數'] = ((ref_date - last_ship['最後出貨日']).dt.days / 30).round(1)
        last_ship['最後出貨日'] = last_ship['最後出貨日'].dt.date
        churn_list = last_ship[last_ship['距今月數'] >= churn_threshold].sort_values('距今月數', ascending=False)

        pct_churn = len(churn_list) / len(last_ship) * 100 if len(last_ship) else 0
        st.metric("流失預警客戶數", f"{len(churn_list)} 家", delta=f"佔全體客戶 {pct_churn:.1f}%")

        if len(churn_list) > 0:
            fig_churn = build_churn_fig(churn_list, churn_threshold)
            st.plotly_chart(fig_churn, use_container_width=True)
            show_table_churn = st.checkbox("顯示表格", value=True, key="t_churn")
            if show_table_churn: st.dataframe(churn_list, use_container_width=True)
        else:
            st.success("目前沒有符合流失警戒門檻的客戶")

    # PPTX 導出：套用畫面上目前的圖表設定重繪，並可選擇套用公司背景範本
    st.sidebar.markdown("### 報表匯出")
    bg_template_file = st.sidebar.file_uploader("上傳背景範本 (.pptx，選填)", type=["pptx"])

    if st.sidebar.button("📊 導出完整戰情室報表"):
        with export_color_theme():
            prs, cfg = make_report_presentation(bg_template_file.getvalue() if bg_template_file else None)

            type_rank_df = f_df.groupby('Machine Type')['Outbound Qty (Item)'].sum().sort_values(ascending=False).reset_index()
            add_dual_chart_slide(
                prs, cfg, "北美設備戰術分佈",
                build_map_fig(f_df), build_type_rank_fig(f_df),
                type_rank_df)

            for name, dim, mode, chart_type, df_g in dim_exports:
                add_chart_slide(prs, cfg, f"分析維度: {name}", build_dim_fig(df_g, dim, mode, chart_type), df_g)

            add_chart_slide(prs, cfg, "客戶 ABC 分析", build_abc_fig(abc_df), abc_df)

            if len(churn_list) > 0:
                add_chart_slide(prs, cfg, f"流失預警（超過 {churn_threshold} 個月無出貨）", build_churn_fig(churn_list, churn_threshold), churn_list)
            else:
                no_churn_slide = prs.slides.add_slide(cfg['layout'])
                if no_churn_slide.shapes.title is not None:
                    no_churn_slide.shapes.title.text = "流失預警"
                else:
                    no_churn_slide.shapes.add_textbox(cfg['left'], Inches(0.2), cfg['width'], Inches(0.5)).text_frame.text = "流失預警"
                no_churn_slide.shapes.add_textbox(cfg['left'], cfg['content_top'], cfg['width'], Inches(0.5)).text_frame.text = "目前沒有符合流失警戒門檻的客戶"

            buf = io.BytesIO()
            prs.save(buf)
        st.sidebar.download_button("下載報表", buf.getvalue(), "Tactical_Report.pptx")
